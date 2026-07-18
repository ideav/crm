#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Сборка PPTX-презентации «Интеграм — для партнёра-реселлера».

Источник контента (истина) — docs/INTEGRAM_PARTNER_DECK_PLAN.md (issue #4273).
Здесь контент курирован под слайды: из плана взяты «Заголовок» (мысль-заголовок),
«Контент»/буллеты, таблицы и картинки; служебные поля («Цель», «Визуал»-описания)
на слайды не выносятся. Номера слайдов 1..24 совпадают с планом (перекрёстные
ссылки «см. слайд N» остаются валидными); слайд 25 — приложение.

Запуск (нужен python-pptx):  python build_partner_deck_pptx.py
Выход: docs/INTEGRAM_PARTNER_DECK.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.normpath(os.path.join(HERE, ".."))
SHOTS = os.path.join(DOCS, "screenshots")
DECK = os.path.join(SHOTS, "deck-4273")
OUT = os.path.join(DOCS, "INTEGRAM_PARTNER_DECK.pptx")

# ---- палитра ----
DARK     = RGBColor(0x0B, 0x2E, 0x4E)   # navy
ACCENT   = RGBColor(0x1E, 0x88, 0xD3)   # синий (как в интерфейсах)
ACCENTDK = RGBColor(0x0E, 0x5A, 0x8F)
TEAL     = RGBColor(0x16, 0x9B, 0x8E)
INK      = RGBColor(0x22, 0x30, 0x3A)
GREY     = RGBColor(0x6A, 0x78, 0x86)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xF2, 0xF6, 0xFA)
LIGHTROW = RGBColor(0xEA, 0xF1, 0xF8)
LACCENT  = RGBColor(0xDC, 0xEC, 0xF7)   # подсветка колонки «Интеграм»
GREEN    = RGBColor(0x1F, 0x9D, 0x55)
FONT     = "Calibri"

EMU_IN = 914400
SW = int(13.333 * EMU_IN)
SH = int(7.5 * EMU_IN)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_counter = 0  # номер слайда для футера


def strip(t):
    return t.replace("`", "").replace("✅", "✓")


def parse_bold(text):
    """'a **b** c' -> [('a ',False),('b',True),(' c',False)]"""
    text = strip(text)
    out, bold = [], False
    for i, part in enumerate(text.split("**")):
        if part:
            out.append((part, i % 2 == 1))
    return out


def _runs(paragraph, text, size, color, bold_all=False, name=FONT):
    for seg, b in parse_bold(text):
        r = paragraph.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.bold = bool(b or bold_all)
        r.font.name = name
        r.font.color.rgb = color


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def footer(slide):
    global _counter
    _counter += 1
    tb, tf = textbox(slide, Inches(0.45), SH - Inches(0.42), Inches(9), Inches(0.3))
    p = tf.paragraphs[0]
    _runs(p, "ИНТЕГРАМ · партнёрская презентация", 9, GREY)
    tb2, tf2 = textbox(slide, SW - Inches(1.2), SH - Inches(0.42), Inches(0.75), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _runs(p2, str(_counter), 9, GREY, bold_all=True)


def title_bar(slide, title, section=None):
    rect(slide, 0, 0, SW, Inches(1.12), DARK)
    rect(slide, 0, Inches(1.12), SW, Inches(0.055), ACCENT)
    tb, tf = textbox(slide, Inches(0.5), 0, SW - Inches(3.4), Inches(1.12), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _runs(p, title, 23, WHITE, bold_all=True)
    if section:
        tb2, tf2 = textbox(slide, SW - Inches(3.2), 0, Inches(2.75), Inches(1.12), MSO_ANCHOR.MIDDLE)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        _runs(p2, section, 11, RGBColor(0xAF, 0xCB, 0xE4))


def subtitle(slide, text, top=Inches(1.32)):
    tb, tf = textbox(slide, Inches(0.5), top, SW - Inches(1.0), Inches(0.5))
    p = tf.paragraphs[0]
    _runs(p, text, 15, ACCENTDK, bold_all=True)
    return top + Inches(0.52)


def _qa_overflow(items, width, height, base):
    """Оценка переполнения текст-бокса (метрики DejaVu ≥ Calibri → консервативно). QA=1."""
    if os.environ.get("QA") != "1":
        return
    try:
        from PIL import ImageFont
    except Exception:
        return
    ttf = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    total = 0.0
    for it in items:
        level, text = it[0], strip(it[1])
        o = it[2] if len(it) > 2 else {}
        size = o.get("size", base - (1.5 if level > 0 else 0))
        f = ImageFont.truetype(ttf, int(round(size)))
        avail = width / EMU_IN * 72 - level * 0.4 * 72 - (0 if o.get("no_bullet") else 22)
        cur, lines = 0.0, 1
        for w in text.split():
            wl = f.getlength(w + " ")
            if cur + wl > avail and cur > 0:
                lines += 1
                cur = wl
            else:
                cur += wl
        total += lines * size * 1.2 * o.get("line_spacing", 1.02) / 72
        total += (o.get("space_after", 5) + o.get("space_before", 0)) / 72
    box = height / EMU_IN
    if total > box + 0.05:
        print(f"[QA overflow] box={box:.2f}in est={total:.2f}in :: {items[0][1][:55]!r}")


def bullets(slide, items, left, top, width, height, base=15.5):
    """items: (level, text, opts?) ; opts: size,color,bold,bullet,bullet_color,no_bullet,space_after"""
    _qa_overflow(items, width, height, base)
    tb, tf = textbox(slide, left, top, width, height)
    first = True
    for it in items:
        level = it[0]
        text = it[1]
        o = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(o.get("space_after", 5))
        p.space_before = Pt(o.get("space_before", 0))
        p.line_spacing = o.get("line_spacing", 1.02)
        size = o.get("size", base - (1.5 if level > 0 else 0))
        if not o.get("no_bullet"):
            indent = "      " * level
            bch = o.get("bullet", "•" if level == 0 else "–")
            r = p.add_run()
            r.text = indent + bch + "  "
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.name = FONT
            r.font.color.rgb = o.get("bullet_color", ACCENT if level == 0 else TEAL)
        _runs(p, text, size, o.get("color", INK), bold_all=o.get("bold", False))
    return tb


def proof(slide, text):
    tb, tf = textbox(slide, Inches(0.5), SH - Inches(0.82), SW - Inches(1.6), Inches(0.42),
                     MSO_ANCHOR.BOTTOM)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Пруф: "
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = TEAL
    _runs(p, text, 10.5, GREY)


def note_box(slide, text, left, top, width, height, fill=LIGHT, color=INK, size=13):
    rect(slide, left, top, width, height, fill)
    tb, tf = textbox(slide, left + Inches(0.15), top, width - Inches(0.3), height, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _runs(p, text, size, color)


def set_cell(cell, text, size, bold, color, fill, align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for seg, b in parse_bold(text):
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.bold = bool(b or bold)
        r.font.name = FONT
        r.font.color.rgb = GREEN if seg.strip().startswith("✓") else color


def table(slide, header, rows, left, top, width, col_widths, font_size=11,
          highlight_last=True, row_h=Inches(0.34), head_h=Inches(0.4)):
    n = len(rows) + 1
    m = len(header)
    gt = slide.shapes.add_table(n, m, left, top, width, head_h + row_h * len(rows)).table
    gt.first_row = False
    gt.horz_banding = False
    for i, w in enumerate(col_widths):
        gt.columns[i].width = w
    gt.rows[0].height = head_h
    for c, h in enumerate(header):
        set_cell(gt.cell(0, c), h, font_size, True, WHITE, DARK,
                 PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for r, row in enumerate(rows, 1):
        gt.rows[r].height = row_h
        for c, val in enumerate(row):
            last = (c == m - 1)
            fill = LACCENT if (highlight_last and last) else (WHITE if r % 2 else LIGHTROW)
            bold = (c == 0) or (highlight_last and last)
            color = ACCENTDK if (highlight_last and last) else INK
            set_cell(gt.cell(r, c), val, font_size, bold, color, fill)
    return gt


def images_row(slide, paths, left, top, width, height, gap=Inches(0.18)):
    pics = []
    for p in paths:
        full = p if os.path.isabs(p) else os.path.join(DECK, p)
        pics.append(slide.shapes.add_picture(full, left, top))
    ars = [pic.width / pic.height for pic in pics]
    tg = gap * (len(pics) - 1)
    h_by_w = (width - tg) / sum(ars)
    h = min(height, h_by_w)
    widths = [h * ar for ar in ars]
    total = sum(widths) + tg
    x = left + (width - total) / 2
    y = top + (height - h) / 2
    for pic, w in zip(pics, widths):
        pic.left = int(x)
        pic.top = int(y)
        pic.width = int(w)
        pic.height = int(h)
        # тонкая рамка
        pic.line.color.rgb = RGBColor(0xCF, 0xDA, 0xE3)
        pic.line.width = Pt(0.75)
        x += w + gap


# =========================================================================
#  СЛАЙД 1 — Титул
# =========================================================================
def slide_title():
    s = add_slide()
    rect(s, 0, 0, SW, SH, DARK)
    rect(s, 0, Inches(4.62), SW, Inches(0.06), ACCENT)
    tb, tf = textbox(s, Inches(0.9), Inches(1.55), SW - Inches(1.8), Inches(1.2))
    p = tf.paragraphs[0]
    _runs(p, "ИНТЕГРАМ", 54, WHITE, bold_all=True)
    tb2, tf2 = textbox(s, Inches(0.9), Inches(2.75), SW - Inches(1.8), Inches(1.0))
    p2 = tf2.paragraphs[0]
    _runs(p2, "low-code платформа на запатентованном ядре данных", 24, RGBColor(0x8F, 0xC3, 0xEC))
    tb3, tf3 = textbox(s, Inches(0.9), Inches(3.75), SW - Inches(2.0), Inches(0.9))
    p3 = tf3.paragraphs[0]
    p3.line_spacing = 1.05
    _runs(p3, "Продукт, который вы продаёте своим заказчикам там, где Excel уже не тянет, "
              "а заказная разработка слишком дорога.", 16, RGBColor(0xD6, 0xE4, 0xF0))
    # регалии
    reg = ("Реестр российского ПО Минцифры №30872   ·   Патенты РФ + US   ·   "
           "Slider-акселератор + финалист SberUp   ·   Продукт недели #1 и #3 Product Radar")
    tb4, tf4 = textbox(s, Inches(0.9), Inches(4.95), SW - Inches(1.8), Inches(0.6))
    p4 = tf4.paragraphs[0]
    _runs(p4, reg, 12.5, RGBColor(0xB9, 0xD3, 0xE8))
    tb5, tf5 = textbox(s, Inches(0.9), Inches(6.15), SW - Inches(1.8), Inches(0.5))
    p5 = tf5.paragraphs[0]
    _runs(p5, "Партнёрская презентация · тон «due diligence продукта» — проверяемые факты, ссылки, цифры.",
          13, RGBColor(0x9A, 0xB4, 0xCC))
    global _counter
    _counter += 1  # титул = слайд 1


# =========================================================================
#  Обобщённый контент-слайд
# =========================================================================
def content(section, title, sub=None, items=None, base=15.5,
            body_left=Inches(0.55), body_top=None, body_w=None, pr=None):
    s = add_slide()
    title_bar(s, title, section)
    top = Inches(1.3)
    if sub:
        top = subtitle(s, sub, top)
    if body_top is not None:
        top = body_top
    if items:
        bullets(s, items, body_left, top, body_w or (SW - Inches(1.1)),
                SH - top - Inches(0.6), base)
    if pr:
        proof(s, pr)
    footer(s)
    return s


def case(section, title, sub, items, imgs, pr=None):
    s = add_slide()
    title_bar(s, title, section)
    top = subtitle(s, sub)
    bullets(s, items, Inches(0.55), top + Inches(0.05), Inches(5.0),
            SH - top - Inches(0.7), base=14.5)
    bottom = (SH - Inches(0.95)) if pr else (SH - Inches(0.55))
    images_row(s, imgs, Inches(5.85), Inches(1.95), SW - Inches(6.35), bottom - Inches(1.95))
    if pr:
        proof(s, pr)
    footer(s)
    return s


# =========================================================================
#  СБОРКА
# =========================================================================
SA = "A · Что это и чем отличается"
SB = "B · Что даёт отличие"
SC = "C · Проблемы · продукт · решения"
SD = "D · Конкуренты · монетизация"

slide_title()  # 1

# --- 2 ---
content(SA, "«Живое ТЗ»: слово клиента сразу становится приложением", items=[
    (0, "От Excel-хаоса до систем корпоративного уровня — без строки кода и без чека заказной разработки."),
    (0, "«Как MS Excel и Access, но **без ограничений** по объёму, скорости, безопасности и интеграциям»."),
    (0, "Даёт: контроль работы людей · автоматизация вычислений · защита данных (ролевая модель) · **не тормозит на любых объёмах**."),
    (0, "Всё проверяется вживую и документируется по ходу — не нужно вычитывать спецификации."),
], base=17)

# --- 3 ---
content(SA, "В основе — квинтетная модель данных (QDM)", "Не реляционка и не голый EAV", items=[
    (0, "Квинтет — атомарная единица: идентификатор · тип · значение · родитель · порядок. Одна таблица вместо сотен, метаданные встроены в данные."),
    (0, "Гибкость EAV **без** его болезни: выборка идёт по индексу (B-деревья), а не сканом → скорость **не падает** с ростом объёма."),
    (0, "Предельная унификация: преимущества РСУБД, NoSQL и колоночных БД; масштаб — распределением по префиксам ID."),
    (0, "Защита: патенты **RU 2650032 C1** и **US 11138174 B2** (одна семья, приоритет 2017, US до 2038). QDM защищён, ядро — открыто."),
    (0, "Конкуренты строят конструкторы поверх реляционки/EAV и упираются в потолок по объёму и логике. Мы — нет.",
        {"color": ACCENTDK, "bold": True, "bullet_color": TEAL}),
], base=15.5, pr="Статья про квинтет — habr.com/ru/companies/neoflex/articles/433058/")

# --- 4 --- матрица «сегмент × заход»
def slide4():
    s = add_slide()
    title_bar(s, "Два захода: заменяем Excel — или становимся ключевой системой", SA)
    bullets(s, [
        (0, "**Направление 1 — замена Excel** (сателлиты): BI, аналитика, планирование, оперативные учётки вокруг основной системы. В корпорациях — основной заход: низкий риск, короткий цикл."),
        (0, "**Направление 2 — замена ключевой системы** (МСБ): производство и глобальное планирование, CRM, ERP. Крупный чек, глубокая интеграция."),
    ], Inches(0.55), Inches(1.35), SW - Inches(1.1), Inches(1.5), base=15)
    table(s,
          ["", "Крупные корпорации", "Малый и средний бизнес (МСБ)"],
          [["Замена Excel\n(BI, аналитика, планирование, сателлиты)",
            "✓ основной заход — вокруг ядра всегда «зоопарк» Excel", "✓ да"],
           ["Замена ключевой системы\n(производство, планирование, CRM, ERP)",
            "реже — есть свой контур", "✓ можем стать core-системой"]],
          Inches(0.55), Inches(3.15), SW - Inches(1.1),
          [Inches(4.2), Inches(4.3), Inches(4.13)], font_size=13,
          highlight_last=False, row_h=Inches(0.95))
    note_box(s, "Оба захода — на одном и том же ядре.", Inches(0.55), Inches(5.75),
             SW - Inches(1.1), Inches(0.6), fill=LACCENT, color=ACCENTDK, size=15)
    footer(s)
slide4()

# --- 5 ---
content(SB, "Ядро открывает то, что конструкторам недоступно", "Несколько ярких примеров — их у нас гораздо больше", items=[
    (0, "Поиск/фильтр по любому полю на **миллиардах записей**", {"bullet": "1.", "size": 16}),
    (0, "Сопоставление номенклатур (matching) без Elasticsearch и разработчиков", {"bullet": "2.", "size": 16}),
    (0, "Память ИИ-агента: вектор + граф + бизнес-данные **в одной БД**", {"bullet": "3.", "size": 16}),
    (0, "Минимальный разрыв с языком бизнеса", {"bullet": "4.", "size": 16}),
    (0, "И это не всё: вычисляемые колонки (LOOKUP/ROLLUP/FORMULA) · шаблонизатор форм · сотни одновременных пользователей · ролевая модель и права · история и откат · автоматизации · вебхуки · импорт батчем (50k за ~30 с) · on-prem · производственное планирование.",
        {"size": 13, "color": GREY, "space_before": 8, "bullet_color": GREY}),
    (0, "Почти каждый такой пример традиционно — отдельный проект на недели/месяцы и команда инженеров.",
        {"color": ACCENTDK, "bold": True, "bullet_color": TEAL, "space_before": 6}),
], base=16)

# --- 6 ---
content(SB, "Пример 1: фильтр по любому полю — на миллиардах записей", "Без денормализации, без Elasticsearch", items=[
    (0, "**32 млрд записей / 4,5+ ТБ** в одной модели (нагрузочный тест), 703+ млн бизнес-транзакций."),
    (0, "Точный поиск по 700 млн записей — **0,77 с**; фильтр по 1 из 45 полей — **1,28 с** на нагруженной базе (~5000 tps)."),
    (0, "Загрузка **~5000 транзакций/с** (32 ядра), без деградации вставки при кратном росте."),
    (0, "Традиционно: денормализация под каждый запрос, Elasticsearch-кластер, шардинг, команда DBA — для no-code классически «нерешаемо».", {"color": GREY}),
    (0, "Контраст: Excel — 1 млн строк, Google Sheets — ~155 тыс. У нас — **безлимит**."),
], pr="habr.com/ru/articles/900308/")

# --- 7 ---
content(SB, "Пример 2: сопоставить чужой каталог с нашей номенклатурой", "За часы, а не недели", items=[
    (0, "Задача: каталог контрагента (22 тыс. позиций) ↔ свои номенклатуры (сотни тысяч) при разных названиях одного товара."),
    (0, "Решение: токенизация названий (regex от ИИ) → справочник токенов → связь «многие-ко-многим» с оценкой качества."),
    (0, "**120–160 сопоставлений/мин, полный каталог за 2–3 часа**, топ-10 кандидатов на позицию; спорное дочищает LLM. Всё на no-code."),
    (0, "Традиционно: ETL + Elasticsearch + отдельный проект разработки. Конструктор такого не умеет.", {"color": GREY}),
    (0, "Уже реализовано и опубликовано — готовый кейс для заказчика."),
], pr="habr.com/ru/articles/1055368/")

# --- 8 ---
content(SB, "Пример 3: память ИИ-агента (VecMory)", "Вектор, граф и бизнес-данные — в одной БД", items=[
    (0, "Обычно это **три хранилища** (вектор-СУБД + граф-СУБД + реляционка) — их надо синхронизировать и обслуживать. У нас — одна модель."),
    (0, "Граф «симптом → причина → фикс»: агент помнит, **на чём уже спотыкался**."),
    (0, "Сбор цепочки причинности — **один серверный вызов ~0,3 с** (против ~33 с клиентского обхода)."),
    (0, "Сублинейный поиск: **×12** при N=5 000 (recall@1 = 1.0), **×86** при N=50 000."),
    (0, "TCO про людей: Pinecone/Qdrant + Neo4j + Postgres = **+0,1–0,3 FTE DevOps/DBA** (≈250–500 тыс ₽/мес). У нас ≈ 0 сверх платформы.", {"color": GREY}),
], base=14.5, pr="vecmory/summary_vm.md в github.com/ideav/crm")

# --- 9 ---
content(SB, "Пример 4: система говорит на языке бизнеса", "Без «Клиент → tbl_cust_main»", items=[
    (0, "Проблема: семантический разрыв «бизнес говорит одно — код называет другое» тормозит внедрение и плодит ошибки."),
    (0, "Решение: сущности — бизнес-термины («Проект», «Задача», «Исполнитель»). Аналитик описывает предметку словами — система строит базу с теми же терминами и связями."),
    (0, "Даёт: быстрый онбординг, меньше ошибок, запросы на бизнес-языке. Задел на будущее — **ИИ работает с бизнес-логикой напрямую**, без слоя перевода."),
    (0, "Традиционно: аналитик → ТЗ → разработчик → БД с техническими именами; каждый переход теряет смысл и деньги.", {"color": GREY}),
], base=16, pr="habr.com/ru/articles/982120/")

# --- 10 ---
def slide10():
    s = add_slide()
    title_bar(s, "Другие конструкторы так не могут — по трём причинам сразу", SB)
    bullets(s, [
        (0, "Хранить **сотни миллионов+** записей."),
        (0, "Реализовывать **сложную логику расчётов без кода**."),
        (0, "Держать **десятки/сотни пользователей** одновременно."),
    ], Inches(0.55), Inches(1.55), SW - Inches(1.1), Inches(2.4), base=20)
    note_box(s, "Это лишь несколько примеров — их гораздо больше. Традиционное решение = кратно больше "
                "ресурсов и сложности; конструктором не решается практически никогда. "
                "Значит, вы берёте проекты, от которых отказываются другие.",
             Inches(0.55), Inches(4.35), SW - Inches(1.1), Inches(1.7), fill=DARK, color=WHITE, size=17)
    footer(s)
slide10()

# --- 11 ---
content(SC, "73% бизнес-задач живут в Excel. Это дорого и опасно.", items=[
    (0, "Слабая дисциплина процессов."),
    (0, "Потери из-за ручного ввода."),
    (0, "Дорогие доработки руками программистов."),
    (0, "Кражи / повреждения данных (в т.ч. ПДн)."),
    (0, "Ограниченность обычного low-code."),
    (0, "McKinsey / BCG / Bain — до **20–30% рабочего времени** уходит на поиск, согласования, рутину.",
        {"color": ACCENTDK, "bold": True, "bullet_color": TEAL, "space_before": 8}),
], base=17)

# --- 12 ---
content(SC, "Любой интерфейс — от таблицы-замены Excel до профессиональной вёрстки", items=[
    (0, "Схематичная вёрстка — быстрый рабочий экран (замена Excel/Airtable)."),
    (0, "Профессиональная вёрстка — клиентский UI."),
    (0, "Адаптив под мобильные."),
    (0, "Формы, отчёты, файлы, права — из коробки."),
], base=17)

# --- 13 --- обзор решений
content(SC, "Решения: пять реальных проектов на Интеграме",
        "Не демо, а работающие системы — от планирования производства до цифрового двойника отрасли",
        items=[
            (0, "**Планирование производства (резка)** — оптимальный план вместо 3 дней ручной работы · ключевая система → слайд 14", {"bullet": "1."}),
            (0, "**Единый дэшборд руководителя** — 13 закладок, автосбор и кросс-валидация · замена Excel → слайд 15", {"bullet": "2."}),
            (0, "**Мобильное РМ оператора (вендинг воды)** — рутина сведена к нулю · ключевая система → слайд 16", {"bullet": "3."}),
            (0, "**Поточный рекрутинг** — 85 тыс. вакансий с HH.ru, 200–300 чел/нед · ключевая система → слайд 17", {"bullet": "4."}),
            (0, "**Дронономика** — цифровой двойник отрасли на онтологиях · уникальный класс → слайд 18", {"bullet": "5."}),
            (0, "Уже сделано и работает — конкретные референсы. Каждый кейс ложится в один из форматов (слайд 20) и одно из двух направлений (слайд 4).",
                {"color": ACCENTDK, "bold": True, "bullet_color": TEAL, "space_before": 8}),
        ], base=16)

# --- 14..18 кейсы ---
case(SC, "Кейс 1: Планирование производства (резка)",
     "3 дня ручного планирования → оптимальный план за один прогон", [
        (0, "**Было:** планирование и перепланирование 5–7 дней производства занимало 3 рабочих дня ручного труда мастера."),
        (0, "**Что сделали:** оптимальный план по множеству факторов с весами — переналадка станков, смена сырья, порядок операций, с учётом физиологии человека и техники станков/сырья."),
        (0, "**Эффект:** план за один прогон и перестройка на лету; освобождены 3 рабочих дня цикла."),
        (0, "**Направление:** ключевая система (операционное ядро производства)."),
     ], ["case-1-planning-a.png", "case-1-planning-b.png"],
     pr="скриншоты — из боевого рабочего места; заказчик под NDA. Код и полный алгоритм планирования открыты — github.com/ideav/crm.")

case(SC, "Кейс 2: Единый дэшборд руководителя",
     "13 разрозненных отчётов → один дэшборд с автосбором и кросс-валидацией", [
        (0, "**Было:** разрозненные материалы на оперативке; дублирующий ручной ввод, нет контроля полноты."),
        (0, "**Что сделали:** 13 закладок (отделы + инвесторские + сводная). Автосбор из Битрикс24 и Google Sheets; загрузка из 1С и Excel; кросс-валидация источников."),
        (0, "**Эффект:** устранён дублирующий ввод; автоматический контроль полноты данных."),
        (0, "**Направление:** замена Excel (BI-сателлит вокруг основной системы)."),
     ], ["case-2-dashboard.png", os.path.join(SHOTS, "connector-1c.png")])

case(SC, "Кейс 3: Мобильное РМ оператора (вендинг воды)",
     "Обслуживание автоматов с телефона — рутина сведена к нулю", [
        (0, "**Было:** копирование данных, ручное планирование задач, сбор отчётов и статистики по аппаратам вручную."),
        (0, "**Что сделали:** мобильное рабочее место для контроля и обслуживания автоматов — задачи, отчёты и статистика в одном экране на телефоне."),
        (0, "**Эффект:** излишние рутинные действия сведены к нулю."),
        (0, "**Направление:** ключевая система (операционное ядро сервисной компании)."),
     ], ["case-3-vending-fleet.png", "case-3-vending-tasks.png", "case-3-vending-report.png"])

case(SC, "Кейс 4: Поточный рекрутинг",
     "Ключевая система найма — за месяц, 200–300 чел/нед без лишних кликов", [
        (0, "**Было:** сложнейшая схема данных, этапов найма и событий; нужна 100% кастомизация под поток."),
        (0, "**Что сделали:** предельная унификация → разработка за 1 месяц. Аналитик правит логику и схему и тут же тестирует — цикл до минут."),
        (0, "**Эффект:** 35 человек одновременно, 85 тыс. вакансий с HH.ru, набор 200–300 чел/нед без лишних кликов."),
        (0, "**Направление:** ключевая система (найм как ядро бизнеса)."),
     ], ["case-4-recruiting-funnel.png", "case-4-recruiting-card.png"])

case(SC, "Кейс 5: Дронономика — цифровой двойник отрасли",
     "Онтологии беспилотной отрасли → живой цифровой двойник с динамическим расчётом", [
        (0, "**Задача:** смоделировать целую отрасль — участников, федеральное управление, экономику, производство — как единую отзывчивую модель."),
        (0, "**Что сделали:** загружены онтологии отрасли, федеральной структуры управления, экономики и производства → точный, отзывчивый цифровой двойник."),
        (0, "**Эффект:** моделирование взаимодействия всех участников с динамическим расчётом параметров — экономика, рейтинги, события, метрики."),
        (0, "**Направление:** уникальный класс (онтологии + граф в одной модели) — мост к ядру и памяти ИИ-агента."),
     ], ["case-5-drone-twin.png"])

# --- 19 безопасность ---
content(SC, "Разворачивается у вас, работает в замкнутом контуре, переживает сбои", items=[
    (0, "Развёртывание и контроль над данными", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15}),
    (1, "Три режима: **SaaS**, **On-premise** (docker на вашем сервере), гибрид."),
    (1, "Замкнутый контур: on-prem полностью изолирован от интернета (**air-gap**) — данные не покидают периметр."),
    (1, "Данные в обычной реляционной СУБД: нет вендор-лока, выгрузка со связями в любой момент."),
    (0, "Безопасность данных", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15, "space_before": 6}),
    (1, "TLS-шифрование каналов + шифрование хранилища."),
    (1, "Ролевая модель **на стороне сервера** — доступ по ролям и маскам."),
    (1, "Корпоративная аутентификация: **JWT, LDAP/AD, SSO**."),
    (0, "Отказоустойчивость", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15, "space_before": 6}),
    (1, "Бэкап + **георезерв между ЦОД**, потоковая репликация; **PITR** на уровне БД; **RPO/RTO под SLA**."),
    (0, "Соответствие: реестр Минцифры №30872 · аттестованное облако под **152-ФЗ и КИИ** · **223-ФЗ / 44-ФЗ**.",
        {"color": ACCENTDK, "bold": True, "bullet_color": TEAL, "space_before": 6}),
], base=14, pr="blog.ideav.ru/posts/bezopasnost-i-otkazoustoichivost-dlya-krupnogo-biznesa/")

# --- 20 форматы ---
content(SC, "Готовые форматы — от Excel-сателлитов до ключевых систем", items=[
    (0, "Замена Excel — сателлиты, короткий цикл", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15}),
    (1, "Тип 1 — учёт договоров, платежей, рабочего времени → **чек до 150 000 ₽**."),
    (1, "BI / аналитика / планирование поверх выгрузок из основной системы."),
    (0, "Ключевая система — МСБ core, средний чек", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15, "space_before": 6}),
    (1, "Тип 2 — контроль и обслуживание аппаратов (вендинг воды) → **чек от 150 000 ₽**."),
    (1, "Тип 3 — рекрутинговый стартап, поток 200–500 чел. → **чек от 500 000 ₽**."),
    (1, "Управление производством и глобальное планирование, CRM, ERP."),
    (0, "Ключевая система — корпорации, крупный чек", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15, "space_before": 6}),
    (1, "Тип 4 — ИТ-сателлиты: трекеры, реестры, BI, регулярные опросы → **чек от 5 млн ₽**."),
    (0, "Шаблоны, которые партнёр/адепт тиражирует; чек растёт со сложностью и переходом «сателлит → ключевая система».",
        {"color": ACCENTDK, "bold": True, "bullet_color": TEAL, "space_before": 6}),
], base=14.5)

# --- 21 конкуренты ---
def slide21():
    s = add_slide()
    title_bar(s, "Мы конкурируем снизу доверху — от Excel до заказной разработки", SD)
    header = ["Критерий", "Excel / Sheets", "Airtable / Coda / Notion", "Bitrix24 / 1С / ELMA", "Заказная разработка", "Интеграм"]
    rows = [
        ["Предел по записям", "1 млн / 155 тыс", "сотни тыс", "зависит", "безлимит", "безлимит (32 млрд проверено)"],
        ["Гибкость логики", "формулы", "ограниченная", "шаблонная", "любая", "любая, no-code"],
        ["Сложная логика без кода", "нет", "ограниченно", "ограниченно", "только кодом", "есть"],
        ["Безопасность / роли", "нет", "базовая", "есть", "есть", "на сервере, по ролям и маскам"],
        ["Локально / замкнутый контур", "локально", "нет", "частично", "да", "on-prem + air-gap"],
        ["Скорость внедрения", "часы", "дни", "недели", "месяцы", "часы–дни"],
        ["Стоимость", "почти бесплатно", "низкая", "средняя", "от сотен тыс", "десятки тыс"],
        ["Импортозамещение / реестр РФ", "—", "нет", "часть", "—", "реестр Минцифры №30872"],
    ]
    cw = [Inches(2.55), Inches(1.7), Inches(2.15), Inches(2.05), Inches(1.95), Inches(1.83)]
    table(s, header, rows, Inches(0.35), Inches(1.35), SW - Inches(0.7), cw,
          font_size=10, row_h=Inches(0.44), head_h=Inches(0.52))
    note_box(s, "Свобода и гибкость Excel без его ограничений. Новые no-code сервисы появляются часто "
                "(МТС Табс, РТ Акола, VK Доска…), но упираются в потолок объёма и логики.",
             Inches(0.35), Inches(6.35), SW - Inches(0.7), Inches(0.62), fill=LACCENT, color=ACCENTDK, size=12.5)
    footer(s)
slide21()

# --- 22 монетизация ---
def slide22():
    s = add_slide()
    title_bar(s, "Партнёрская модель: где здесь партнёр и как он зарабатывает", SD)
    subtitle(s, "Партнёр даёт канал и одобряет сделки; внедряют и держат поддержку адепты — доход без операционной нагрузки")
    table(s,
          ["Роль", "Что делает", "На чём зарабатывает"],
          [["Платформа (Интеграм)", "развивает продукт, обучает адептов, выдаёт лицензии", "подписка, лицензия, 20–50% роялти с тиражных решений"],
           ["Партнёр (вы)", "инициирует и одобряет сделки, приводит заказчиков и бренд; не внедряет", "% с каждой оплаты клиента (40/20/15%) — без операционной нагрузки"],
           ["Адепт", "собирает решение, внедряет, ведёт поддержку, отвечает за результат", "оплата работы — внутри чека проекта"],
           ["Заказчик", "конечный клиент партнёра", "платит за проект и поддержку; платформе — за подписку"]],
          Inches(0.4), Inches(1.95), Inches(8.1),
          [Inches(1.7), Inches(3.2), Inches(3.2)], font_size=11,
          highlight_last=False, row_h=Inches(0.7), head_h=Inches(0.36))
    # ставки
    table(s,
          ["Тип продажи", "Партнёру"],
          [["Локальная лицензия Интеграм", "40 %"],
           ["Проект автоматизации", "20 %"],
           ["SaaS-платежи / поддержка", "15 %"]],
          Inches(8.75), Inches(1.95), Inches(4.15),
          [Inches(3.0), Inches(1.15)], font_size=12,
          highlight_last=False, row_h=Inches(0.5), head_h=Inches(0.4))
    note_box(s, "Плата платформе (подписка от 1 950 ₽/мес или лицензия 590 тыс ₽/год) — стоимость инструмента, "
                "а НЕ комиссия с проекта. Два потока дохода партнёра: разовая доля с проектов + рекуррент "
                "с поддержки, доработок и лицензий.",
             Inches(8.75), Inches(3.85), Inches(4.15), Inches(1.6), fill=LIGHT, color=INK, size=11.5)
    note_box(s, "Платформа обучает адептов  →  Партнёр приводит и одобряет сделку  →  Адепт внедряет и поддерживает  →  "
                "Заказчик платит; партнёр берёт свою долю, не работая руками.",
             Inches(0.4), Inches(5.95), Inches(8.1), Inches(0.92), fill=DARK, color=WHITE, size=12.5)
    footer(s)
slide22()

# --- 23 финмодель ---
content(SD, "2025–26 — фокус на инструменте и стратегическом партнёрстве", items=[
    (0, "**Трекшн:** выручка 2022 → 2023 → 2024 = **0,2 → 0,7 → 3,1 млн ₽** (×4–5/год); 70 активных клиентов РФ/СНГ; 5 лицензий; 5 партнёров + 7 адептов зарабатывают; +300 адептов обучены; финалист SberUp; Продукт недели #1 Product Radar."),
    (0, "Фокус 2025–26", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15, "space_before": 4}),
    (1, "Оттачивание инструмента и стратегическое партнёрство — не гонка за выручкой."),
    (1, "**5 активных партнёров** уже зарабатывают; **ещё 10** знакомят своих заказчиков."),
    (1, "КП поданы в **ВТБ, Газпром нефть (ГПН)** и др.; референс-проекты + адепты — фундамент экосистемы."),
    (1, "Ноябрь 2025 — в **Реестре ПО Минцифры** (№30872); открытый стек, ИИ-интеграция, on-prem."),
    (0, "**Рынок:** TAM **4,5 трлн ₽** ($44,5B) · SAM **1,8 трлн ₽** ($18B) · SOM **4 млрд ₽** ($40M); CAGR **19,2%**; мировой low-code → **$77 млрд к 2030**.", {"space_before": 4}),
    (0, "Отдела продаж нет и не будет — приходят и платят органически. Партнёр делает спрос системным.",
        {"color": ACCENTDK, "bold": True, "bullet_color": TEAL}),
], base=13.5)

# --- 24 команда/оффер/контакты ---
def slide24():
    s = add_slide()
    title_bar(s, "Опытная команда, проверяемая технология, готовы дать доступ", SD)
    bullets(s, [
        (0, "**Команда:** CEO Алексей Семёнов (разработчик ядра; Citi, MTS, Neoflex, ВТБ, Сбер), CTO Александр Орехов (highload, 15 лет), аналитики Ким Артамонов и Денис Гаврилов; +300 адептов."),
        (0, "**Пруфы:** патенты (Google Patents RU/US), боевой репозиторий github.com/ideav/crm (3700+ тикетов), публикации на Хабре, реестр Минцифры, SberUp."),
        (0, "Оффер партнёру", {"bold": True, "no_bullet": True, "color": ACCENTDK, "size": 15, "space_before": 4}),
        (1, "Продукт в портфель без затрат на разработку — на готовом запатентованном ядре."),
        (1, "Обучение и сертификация адептов; вы даёте канал и одобряете сделки. Доля с проектов + рекуррент с подписок/поддержки."),
        (1, "Готовые форматы и референс-проекты; маркетинговая и техническая поддержка; гибкие модели (реселлер / внедренческий / стратегический / эксклюзивный)."),
        (0, "**Next steps:** 1) демо + доступ к боевым проектам  2) пилот на одном заказчике  3) обучение первых адептов  4) партнёрское соглашение.", {"space_before": 4}),
    ], Inches(0.55), Inches(1.35), SW - Inches(1.1), Inches(4.0), base=13.5)
    note_box(s, "Алексей Семёнов, CEO   ·   Telegram @qdmadept   ·   abc@integram.io   ·   +7 (995) 506-01-67   ·   ideav.ru   ·   help.integram.io",
             Inches(0.55), Inches(5.55), SW - Inches(1.1), Inches(0.6), fill=LIGHT, color=INK, size=12.5)
    note_box(s, "Избавим бизнес от рутины и защитим данные. «Tetra Pak в сфере IT».",
             Inches(0.55), Inches(6.25), SW - Inches(1.1), Inches(0.62), fill=DARK, color=WHITE, size=15)
    footer(s)
slide24()

# --- 25 приложение ---
content(SD, "Приложение (по запросу заказчика / партнёра)", items=[
    (0, "Технические замеры VecMory — сравнение с pgvector / Neo4j, TCO-таблица."),
    (0, "Разбор кейсов подробнее — скриншоты, ТЗ, чек."),
    (0, "Безопасность и отказоустойчивость — детально (SaaS / on-prem / air-gap, RPO/RTO, 152-ФЗ / КИИ)."),
    (0, "Патентные документы, выписка из реестра Минцифры."),
    (0, "Прайс и партнёрский договор."),
], base=16)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
